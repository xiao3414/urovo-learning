"""从桌面课件目录提取 PPT 文本（不复制 PPT 到项目中）

用法:
  pip install python-pptx
  python scripts/extract_pptx.py

输出: scripts/pptx_extracted.json（仅供开发参考，需手动整理进 src/data/trainingContent.js）
PPT 原文件保留在桌面，不进入 Git 仓库。
"""
import json
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
DESKTOP = Path(r'c:\Users\urovo\Desktop\集训课件')

SRC_MAP = {
    'day1-2-basic': DESKTOP / '1-集训课件day1&day2' / '1-集训课件day1&day2' / '基础认识篇',
    'day1-2-business': DESKTOP / '1-集训课件day1&day2' / '1-集训课件day1&day2' / '深入业务篇',
    'day3-4-business': DESKTOP / '2-集训课件day3&day4' / '2-集训课件day3&day4' / '深入业务篇',
}


def extract_pptx(path):
    from pptx import Presentation
    texts = []
    prs = Presentation(path)
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            texts.append('\n'.join(parts))
    return texts


if __name__ == '__main__':
    all_slides = {}
    for folder, src in SRC_MAP.items():
        if not src.is_dir():
            print('SKIP (目录不存在):', folder, '->', src)
            continue
        for fn in sorted(src.glob('*.pptx')):
            if fn.name.startswith('~$'):
                continue
            try:
                all_slides[fn.stem] = extract_pptx(fn)
                print('OK', folder, fn.name, len(all_slides[fn.stem]), 'slides')
            except Exception as e:
                all_slides[fn.stem] = [f'ERROR: {e}']
                print('ERR', fn.name, e)

    out = ROOT / 'scripts' / 'pptx_extracted.json'
    with open(out, 'w', encoding='utf-8') as f:
        json.dump(all_slides, f, ensure_ascii=False, indent=2)
    print(f'Done -> {out}')
    print('下一步: 将要点整理进 src/data/trainingContent.js，勿提交 PPT 或大体积 JSON。')
