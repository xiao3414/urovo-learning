"""Extract text from WXWork training files for review."""
import json
from pathlib import Path

BASE = Path(r'c:\Users\urovo\Documents\WXWork\1688854310865517\Cache\File\2026-08')
OUT = Path(__file__).resolve().parent / 'wxwork_extracted.json'

FILES = [
    '1、AIDC基础知识及行业概况.pdf',
    '2、优博讯通用版本PDA全系类产品培训.pdf',
    '3、通用产品研发流程培训.pptx',
    '优博讯全系列产品介绍——中文20260131 (1).pptx',
]


def extract_pdf(path: Path) -> list[str]:
    try:
        import fitz
        doc = fitz.open(path)
        pages = []
        for page in doc:
            text = page.get_text().strip()
            if text:
                pages.append(text)
        return pages
    except Exception:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        return [(p.extract_text() or '').strip() for p in reader.pages if (p.extract_text() or '').strip()]


def extract_pptx(path: Path) -> list[str]:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            slides.append('\n'.join(parts))
    return slides


if __name__ == '__main__':
    result = {}
    for name in FILES:
        path = BASE / name
        if not path.exists():
            print('MISSING', name)
            continue
        ext = path.suffix.lower()
        try:
            if ext == '.pdf':
                result[name] = extract_pdf(path)
            elif ext == '.pptx':
                result[name] = extract_pptx(path)
            print('OK', name, len(result[name]), 'pages/slides')
        except Exception as e:
            result[name] = [f'ERROR: {e}']
            print('ERR', name, e)

    OUT.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding='utf-8')
    print('Saved', OUT)
